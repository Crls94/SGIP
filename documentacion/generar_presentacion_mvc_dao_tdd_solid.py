from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = "/home/chris/Escritorio/Metro - Proy - Integrador/sgip-backend/documentacion/presentacion_mvc_dao_tdd_solid_sgip.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "navy": RGBColor(18, 47, 74),
    "blue": RGBColor(31, 78, 121),
    "light_blue": RGBColor(221, 235, 247),
    "green": RGBColor(84, 130, 53),
    "orange": RGBColor(197, 90, 17),
    "gray": RGBColor(89, 89, 89),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(25, 25, 25),
}


def set_text(frame, text, size=24, color="black", bold=False, align=None):
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Calibri"
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    if align:
        p.alignment = align


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.65))
    set_text(box.text_frame, title, 30, "navy", True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.fill.background()
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.08), Inches(12), Inches(0.35))
        set_text(sub.text_frame, subtitle, 13, "gray")


def add_footer(slide, number):
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"SGIP Metro - MVC, DAO, TDD y SOLID | {number}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["gray"]
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items, x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS["black"]
        p.space_after = Pt(8)
    return box


def add_card(slide, x, y, w, h, title, body, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light_blue"]
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(1.4)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS[color]
    p.space_after = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLORS["black"]
    return shape


def add_table_like(slide, headers, rows, x, y, w, h, font_size=12):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = COLORS["white"]
            p.font.bold = True
            p.font.size = Pt(font_size)
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["black"]
    return table


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = COLORS["navy"]
title = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.8), Inches(1.2))
set_text(title.text_frame, "MVC, DAO, TDD, SOLID y CRUD", 42, "white", True, PP_ALIGN.CENTER)
subtitle = slide.shapes.add_textbox(Inches(1.3), Inches(2.85), Inches(10.8), Inches(0.9))
set_text(subtitle.text_frame, "Aplicado al proyecto SGIP - Metro Ica", 26, "white", False, PP_ALIGN.CENTER)
body = slide.shapes.add_textbox(Inches(1.25), Inches(4.25), Inches(10.8), Inches(0.8))
set_text(body.text_frame, "Sistema de Gestión Inteligente de Inventarios y Pedidos", 20, "white", False, PP_ALIGN.CENTER)
add_footer(slide, 1)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Objetivo de la arquitectura", "Por qué se organiza así el sistema SGIP")
add_bullets(slide, [
    "Separar responsabilidades para que una clase no haga todo.",
    "Mantener módulos independientes: productos, pedidos, movimientos, seguridad, alertas.",
    "Facilitar pruebas unitarias con JUnit antes o durante la implementación.",
    "Conectar frontend React, backend Spring Boot y PostgreSQL mediante capas claras.",
], 0.9, 1.65, 11.5, 4.1, 22)
add_footer(slide, 2)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Terminologías clave", "Resumen breve para entender la propuesta")
add_card(slide, 0.65, 1.45, 3.8, 1.35, "MVC", "Model, View, Controller: separa datos, interfaz y controladores.")
add_card(slide, 4.75, 1.45, 3.8, 1.35, "DAO", "Objeto de acceso a datos. En SGIP son los Repository de JPA.", "green")
add_card(slide, 8.85, 1.45, 3.8, 1.35, "DTO", "Transporta datos entre capas sin exponer entidades completas.", "orange")
add_card(slide, 0.65, 3.35, 3.8, 1.35, "TDD", "Primero prueba, luego código. Reduce errores de negocio.", "green")
add_card(slide, 4.75, 3.35, 3.8, 1.35, "SOLID", "Principios para código limpio, mantenible y desacoplado.")
add_card(slide, 8.85, 3.35, 3.8, 1.35, "CRUD", "Crear, leer, actualizar y eliminar datos.", "orange")
add_footer(slide, 3)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 1: ¿Qué elementos forman MVC?", "Aplicado al proyecto real SGIP Metro")
add_table_like(slide, ["Capa", "Qué representa", "Ejemplos reales"], [
    ["Modelo", "Datos y reglas del negocio", "Producto.java, Pedido.java, Usuario.java, InventarioMovimiento.java"],
    ["Vista", "Interfaz que usa el usuario", "Frontend React: Login, Dashboard, Productos, Pedidos, Alertas"],
    ["Controlador", "Recibe peticiones HTTP y llama al Service", "ProductoController.java, PedidoController.java, MovimientoController.java"],
], 0.55, 1.55, 12.2, 3.2, 14)
add_bullets(slide, ["Flujo ejemplo: React -> ProductoController -> ProductoService -> ProductoRepository -> PostgreSQL"], 0.85, 5.2, 11.7, 0.8, 20)
add_footer(slide, 4)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 1: Datos y acceso DAO", "Qué se almacena y cómo se accede a la base de datos")
add_table_like(slide, ["Datos", "Modelo", "DAO / Repository"], [
    ["Productos", "Producto", "ProductoRepository"],
    ["Pedidos", "Pedido / PedidoDetalle", "PedidoRepository / PedidoDetalleRepository"],
    ["Inventario", "InventarioMovimiento", "MovimientoRepository"],
    ["Usuarios", "Usuario / Sesion", "UsuarioRepository / SesionRepository"],
    ["Alertas", "AlertaStock", "AlertaStockRepository"],
], 0.55, 1.4, 12.2, 4.1, 13)
add_bullets(slide, [
    "Los DAO están implementados con Spring Data JPA.",
    "Ejemplo: ProductoRepository extends JpaRepository<Producto, UUID>.",
], 0.8, 5.8, 11.8, 0.75, 18)
add_footer(slide, 5)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 1: Prueba TDD elegida", "Primero se define la prueba, después la lógica")
add_card(slide, 0.85, 1.55, 11.65, 1.3, "Prueba en una frase", "El sistema debe impedir una salida de inventario cuando la cantidad solicitada sea mayor al stock actual del producto.", "orange")
add_bullets(slide, [
    "Dado un producto con stock actual de 5 unidades.",
    "Cuando un usuario intenta registrar una salida de 10 unidades.",
    "Entonces el sistema debe rechazar la operación con el mensaje: Stock insuficiente.",
    "Esta regla corresponde a MovimientoService.registrarMovimientoReal().",
], 1.0, 3.25, 11.3, 2.3, 21)
add_footer(slide, 6)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 1: Principio SOLID principal", "Cómo evitar que una clase haga demasiadas cosas")
add_card(slide, 0.85, 1.35, 11.65, 0.95, "SRP - Single Responsibility Principle", "Cada clase debe tener una sola responsabilidad y una sola razón para cambiar.", "green")
add_table_like(slide, ["Clase", "Responsabilidad"], [
    ["ProductoController", "Recibir peticiones HTTP"],
    ["ProductoService", "Aplicar reglas de negocio"],
    ["ProductoRepository", "Acceder a PostgreSQL"],
    ["ProductoMapper", "Convertir entidades y DTOs"],
], 1.1, 2.65, 11.1, 3.2, 14)
add_footer(slide, 7)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 2: Diseño rápido de capas", "Diagrama simple MVC + DAO")
steps = [
    ("Vista", "Frontend React\nProductos, Pedidos, Dashboard"),
    ("Controlador", "Spring Boot REST\nProductoController, PedidoController"),
    ("Modelo + Lógica", "Entidades, DTOs, Services\nProducto, Pedido, ProductoService"),
    ("DAO / Repository", "Spring Data JPA\nProductoRepository, PedidoRepository"),
    ("Base de datos", "PostgreSQL"),
]
y = 1.35
for i, (head, body_text) in enumerate(steps):
    add_card(slide, 3.2, y + i * 1.05, 6.9, 0.78, head, body_text, "blue" if i < 3 else "green")
    if i < len(steps) - 1:
        arrow = slide.shapes.add_textbox(Inches(6.25), Inches(y + i * 1.05 + 0.72), Inches(0.7), Inches(0.25))
        set_text(arrow.text_frame, "↓", 18, "gray", True, PP_ALIGN.CENTER)
add_footer(slide, 8)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Parte 2: DAO, TDD y SOLID en el diseño", "Respuesta final del grupo")
add_card(slide, 0.7, 1.3, 5.9, 1.35, "Dónde van los DAO", "Dentro de cada módulo: productos/ProductoRepository, pedidos/PedidoRepository, movimientos/MovimientoRepository, seguridad/UsuarioRepository.")
add_card(slide, 6.85, 1.3, 5.75, 1.35, "Prueba TDD", "Impedir una salida de inventario cuando la cantidad solicitada sea mayor al stock actual.", "orange")
add_table_like(slide, ["Principio", "Cómo se respeta"], [
    ["SRP", "Controller, Service, Repository, Mapper y DTO tienen tareas separadas."],
    ["DIP", "Los Services dependen de interfaces Repository, no de SQL directo."],
    ["OCP", "Se pueden agregar nuevos módulos o reglas sin reescribir todo."],
], 0.8, 3.25, 11.8, 2.55, 14)
add_footer(slide, 9)

prs.save(OUTPUT)
print(OUTPUT)
